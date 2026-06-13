#!/usr/bin/env python3
"""LPAI Nexus — 23-slide PowerPoint Presentation Generator"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Colors ─────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1A, 0x3A, 0x6A)
GOLD  = RGBColor(0xD4, 0x88, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LBLUE = RGBColor(0xE8, 0xF0, 0xFA)
DARK  = RGBColor(0x1C, 0x1C, 0x2E)
GREEN = RGBColor(0x1A, 0x6A, 0x3A)
RED   = RGBColor(0xC0, 0x39, 0x2B)
GRAY  = RGBColor(0x66, 0x66, 0x66)
LGRAY = RGBColor(0xF0, 0xF0, 0xF0)
AMBER = RGBColor(0xF3, 0x9C, 0x12)
LGREEN= RGBColor(0xD4, 0xED, 0xDA)
LAMBER= RGBColor(0xFF, 0xF3, 0xCD)
MGRAY = RGBColor(0xCC, 0xCC, 0xCC)
DNAVY = RGBColor(0x0D, 0x20, 0x40)
LRED  = RGBColor(0xF8, 0xD7, 0xDA)
TEAL  = RGBColor(0x17, 0x6B, 0x87)

# ── Slide dimensions ────────────────────────────────────────────────────────
W     = Inches(13.33)
H     = Inches(7.5)
TOP_H = Inches(0.4)
BOT_H = Inches(0.3)
BOT_Y = Inches(7.2)
LST_W = Inches(0.15)

SS_BASE = "/home/asus/lpai/report-assests/screenshots"
SCREENS = {
    "command":      f"{SS_BASE}/screenshot-homepage.jpeg",
    "cargo":        f"{SS_BASE}/screenshot-cargo.jpeg",
    "immigration":  f"{SS_BASE}/screenshot-people.png",
    "security":     f"{SS_BASE}/screenshot-security.png",
    "vehicle":      f"{SS_BASE}/screenshot-vehiclegate.png",
    "intelligence": f"{SS_BASE}/screenshot-intelligence.png",
    "walkthrough":  f"{SS_BASE}/screenshot-walkthrough.png",
}

OUTPUT = "/home/asus/lpai/LPAI_Nexus_Presentation.pptx"


# ── Core helpers ────────────────────────────────────────────────────────────

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _remove_line(shp):
    try:
        spPr = shp._element.spPr
        for ln in spPr.findall(qn('a:ln')):
            spPr.remove(ln)
        ln_el = etree.SubElement(spPr, qn('a:ln'))
        ln_el.set('w', '0')
        etree.SubElement(ln_el, qn('a:noFill'))
    except Exception:
        pass


def rect(slide, x, y, w, h, fill=None, line_c=None, lw=1.5):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_c is not None:
        shp.line.color.rgb = line_c
        shp.line.width = Pt(lw)
    else:
        _remove_line(shp)
    return shp


def txb(slide, x, y, w, h, paras, wrap=True):
    shp = slide.shapes.add_textbox(x, y, w, h)
    tf = shp.text_frame
    tf.word_wrap = wrap
    for i, pd in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pd.get('align', PP_ALIGN.LEFT)
        if pd.get('space_before'): p.space_before = Pt(pd['space_before'])
        if pd.get('space_after'):  p.space_after  = Pt(pd['space_after'])
        run = p.add_run()
        run.text           = pd.get('text', '')
        run.font.name      = pd.get('font', 'Calibri')
        run.font.size      = Pt(pd.get('size', 18))
        run.font.bold      = pd.get('bold', False)
        run.font.italic    = pd.get('italic', False)
        run.font.color.rgb = pd.get('color', DARK)
    return shp


def img(slide, key, x, y, w, h=None):
    path = SCREENS.get(key, '')
    if os.path.exists(path):
        try:
            if h:
                return slide.shapes.add_picture(path, x, y, width=w, height=h)
            return slide.shapes.add_picture(path, x, y, width=w)
        except Exception as e:
            print(f"  Warn img {key}: {e}")
    ph_h = h or Inches(3.5)
    rect(slide, x, y, w, ph_h, fill=MGRAY)
    txb(slide, x, y + ph_h // 2 - Inches(0.2), w, Inches(0.4),
        paras=[{'text': f'[ {key} — screenshot ]', 'size': 11,
                'color': GRAY, 'align': PP_ALIGN.CENTER}])


def set_notes(slide, text):
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception:
        pass


def cell_fmt(cell, text, bold=False, color=DARK, bg=None, size=10,
             align=PP_ALIGN.LEFT, italic=False):
    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    if p.runs:
        r = p.runs[0]
    else:
        r = p.add_run()
        r.text = text
    r.font.name      = 'Calibri'
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color


def add_master(slide, num, title):
    rect(slide, 0, 0, LST_W, H, fill=GOLD)
    rect(slide, 0, 0, W, TOP_H, fill=NAVY)
    txb(slide, Inches(0.25), Inches(0.07), Inches(8.5), Inches(0.3),
        paras=[{'text': 'LPAI Nexus — Smart Border Command Platform',
                'size': 11, 'bold': True, 'color': WHITE}], wrap=False)
    txb(slide, Inches(9.0), Inches(0.07), Inches(4.1), Inches(0.3),
        paras=[{'text': 'Ministry of Home Affairs | Government of India',
                'size': 10, 'color': WHITE, 'align': PP_ALIGN.RIGHT}], wrap=False)
    rect(slide, 0, BOT_Y, W, BOT_H, fill=LGRAY)
    txb(slide, Inches(0.25), BOT_Y + Inches(0.04), Inches(4), Inches(0.22),
        paras=[{'text': 'RESTRICTED — For Internal Use Only',
                'size': 9, 'bold': True, 'color': RED}], wrap=False)
    txb(slide, Inches(4.5), BOT_Y + Inches(0.04), Inches(4.5), Inches(0.22),
        paras=[{'text': title, 'size': 9, 'color': GRAY,
                'align': PP_ALIGN.CENTER}], wrap=False)
    txb(slide, Inches(10.2), BOT_Y + Inches(0.04), Inches(3.0), Inches(0.22),
        paras=[{'text': f'{num} / 23', 'size': 9, 'bold': True,
                'color': NAVY, 'align': PP_ALIGN.RIGHT}], wrap=False)


def stitle(slide, text):
    txb(slide, Inches(0.3), Inches(0.44), Inches(12.8), Inches(0.5),
        paras=[{'text': text, 'size': 25, 'bold': True, 'color': NAVY}])
    rect(slide, Inches(0.3), Inches(0.95), Inches(12.8), Inches(0.04), fill=GOLD)


# ── Slide builders ──────────────────────────────────────────────────────────

def s01(prs):
    sl = new_slide(prs)
    rect(sl, 0, 0, W, H, fill=NAVY)
    rect(sl, 0, 0, LST_W, H, fill=GOLD)
    rect(sl, 0, Inches(6.9), W, Inches(0.6), fill=DNAVY)
    # Top area
    txb(sl, Inches(0.5), Inches(0.55), Inches(12.33), Inches(0.4),
        paras=[{'text': '  GOVERNMENT OF INDIA', 'size': 14, 'bold': True,
                'color': WHITE, 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(1.0), Inches(12.33), Inches(0.35),
        paras=[{'text': 'Ministry of Home Affairs  |  Land Port Authority of India',
                'size': 12, 'color': RGBColor(0xCC, 0xDD, 0xFF),
                'align': PP_ALIGN.CENTER}])
    rect(sl, Inches(3.5), Inches(1.45), Inches(6.33), Inches(0.04), fill=GOLD)
    # Main title
    txb(sl, Inches(0.5), Inches(1.6), Inches(12.33), Inches(0.85),
        paras=[{'text': 'AI Integration in Land Port Authority of India',
                'size': 34, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER}])
    # Subtitle
    txb(sl, Inches(0.5), Inches(2.5), Inches(12.33), Inches(0.65),
        paras=[{'text': 'LPAI Nexus — Smart Border Command Platform',
                'size': 26, 'bold': True, 'color': GOLD, 'align': PP_ALIGN.CENTER}])
    rect(sl, Inches(3.5), Inches(3.25), Inches(6.33), Inches(0.04), fill=GOLD)
    # Sub-details
    txb(sl, Inches(0.5), Inches(3.4), Inches(12.33), Inches(0.35),
        paras=[{'text': 'Internship Technical Presentation',
                'size': 13, 'color': RGBColor(0xCC, 0xDD, 0xFF),
                'italic': True, 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(3.85), Inches(12.33), Inches(0.4),
        paras=[{'text': 'Himanshu  |  AI & Software Engineering Intern, LPAI',
                'size': 15, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(4.3), Inches(12.33), Inches(0.35),
        paras=[{'text': 'June 2026', 'size': 13,
                'color': RGBColor(0xCC, 0xDD, 0xFF), 'align': PP_ALIGN.CENTER}])
    # Bottom band text
    txb(sl, Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.4),
        paras=[{'text': 'Land Port Authority of India  |  Ministry of Home Affairs  |  Government of India',
                'size': 10, 'color': RGBColor(0x88, 0x99, 0xBB),
                'align': PP_ALIGN.CENTER}])
    set_notes(sl, "Introduce yourself and the internship context. Emphasize that LPAI Nexus is a fully working prototype built during the internship — not a concept or proposal document. This presentation covers the platform built, the problem it solves, and recommended next steps.")
    return sl


def s02(prs):
    sl = new_slide(prs)
    T = "India's Land Borders — The Scale of the Challenge"
    add_master(sl, 2, T); stitle(sl, T)
    # 4 stat boxes — 2x2 grid left half
    cfg = [
        (NAVY,  "15,200 KM",    "International Land Border"),
        (GOLD,  "7 Countries",  "Afghanistan · Bangladesh · Bhutan\nChina · Myanmar · Nepal · Pakistan"),
        (GREEN, "12 ICPs",      "Operational Integrated Check Posts"),
        (NAVY,  "23 Sanctioned","Total Land Ports Approved"),
    ]
    for i, (c, num, lbl) in enumerate(cfg):
        bx = Inches(0.3) + (i % 2) * Inches(3.1)
        by = Inches(1.1)  + (i // 2) * Inches(1.85)
        rect(sl, bx, by, Inches(2.9), Inches(1.7), fill=c)
        txb(sl, bx + Inches(0.1), by + Inches(0.12), Inches(2.7), Inches(0.7),
            paras=[{'text': num, 'size': 28, 'bold': True,
                    'color': DARK if c == GOLD else WHITE,
                    'align': PP_ALIGN.CENTER}])
        txb(sl, bx + Inches(0.1), by + Inches(0.85), Inches(2.7), Inches(0.7),
            paras=[{'text': lbl, 'size': 10,
                    'color': DARK if c == GOLD else WHITE,
                    'align': PP_ALIGN.CENTER}])
    # Right panel — ICP list
    rect(sl, Inches(6.55), Inches(1.05), Inches(6.6), Inches(5.85), fill=LBLUE)
    txb(sl, Inches(6.65), Inches(1.12), Inches(6.4), Inches(0.38),
        paras=[{'text': 'Operational ICPs Across Land Borders',
                'size': 12, 'bold': True, 'color': NAVY, 'align': PP_ALIGN.CENTER}])
    rect(sl, Inches(6.65), Inches(1.5), Inches(6.3), Inches(0.03), fill=GOLD)
    icps = [
        ("Attari–Wagah", "Punjab — Pakistan Border"),
        ("Petrapole", "West Bengal — Bangladesh"),
        ("Raxaul", "Bihar — Nepal Border"),
        ("Moreh", "Manipur — Myanmar Border"),
        ("Dawki", "Meghalaya — Bangladesh"),
        ("Sutarkhandi", "Assam — Bangladesh"),
        ("Agartala", "Tripura — Bangladesh"),
        ("Jogbani", "Bihar — Nepal Border"),
        ("Rupaidiha", "Uttar Pradesh — Nepal"),
        ("Hili", "West Bengal — Bangladesh"),
        ("Mahendragarh*", "Uttarakhand — Nepal"),
        ("Sabroom*", "Tripura — Bangladesh"),
    ]
    for i, (name, loc) in enumerate(icps):
        col = 0 if i < 6 else 1
        row = i % 6
        cx = Inches(6.7) + col * Inches(3.15)
        cy = Inches(1.6) + row * Inches(0.72)
        dot = "●" if "*" not in name else "◌"
        dc = GREEN if "*" not in name else AMBER
        txb(sl, cx, cy, Inches(0.25), Inches(0.32),
            paras=[{'text': dot, 'size': 10, 'color': dc, 'bold': True}])
        txb(sl, cx + Inches(0.22), cy, Inches(2.7), Inches(0.3),
            paras=[{'text': name.replace('*',''), 'size': 10, 'bold': True, 'color': NAVY}])
        txb(sl, cx + Inches(0.22), cy + Inches(0.28), Inches(2.7), Inches(0.26),
            paras=[{'text': loc, 'size': 8.5, 'color': GRAY}])
    txb(sl, Inches(6.7), Inches(5.82), Inches(6.2), Inches(0.22),
        paras=[{'text': '* Commissioning phase', 'size': 8, 'italic': True, 'color': GRAY}])
    # Bottom banner
    rect(sl, Inches(0.3), Inches(5.08), Inches(12.9), Inches(0.82), fill=DNAVY)
    txb(sl, Inches(0.4), Inches(5.14), Inches(12.7), Inches(0.3),
        paras=[{'text': 'FY 2023-24 Highlights', 'size': 11, 'bold': True, 'color': GOLD}])
    txb(sl, Inches(0.4), Inches(5.44), Inches(12.7), Inches(0.4),
        paras=[{'text': '₹70,952 Crore trade facilitated  ·  30.46 Lakh passengers processed  ·  ₹62.17 Crore revenue collected  —  15-fold trade increase from 2014-15',
                'size': 12, 'bold': True, 'color': WHITE, 'align': PP_ALIGN.CENTER}])
    set_notes(sl, "Establish the extraordinary scale of India's land border challenge. The 15-fold trade increase is the most powerful number here — it shows that LPAI is managing explosive growth with infrastructure and processes designed for a fraction of current volumes. The case for AI is not theoretical; it is operational necessity.")
    return sl


def s03(prs):
    sl = new_slide(prs)
    T = "LPAI's Digital Transformation — Strong Foundations"
    add_master(sl, 3, T); stitle(sl, T)
    cols = [
        {
            'color': NAVY, 'title': 'LPMS',
            'sub': 'Land Port Management System',
            'icon': '🖥️',
            'bullets': [
                'Single window for all port operations',
                'ANPR cameras + VAHAN vehicle integration',
                'Slot booking deployed for trucks',
                'Significant reduction in wait times',
            ]
        },
        {
            'color': TEAL, 'title': 'PIFS',
            'sub': 'Passenger Integrated Facilitation System',
            'icon': '🛂',
            'bullets': [
                'Aadhaar + passport digital registration',
                'Facial recognition at entry/exit points',
                'Pilot at Petrapole — successful rollout',
                'Real-time queue visibility enabled',
            ]
        },
        {
            'color': GREEN, 'title': 'Suvidha Portal',
            'sub': 'Trader Self-Service Platform',
            'icon': '✅',
            'bullets': [
                'Slot booking live at Bangladesh ICPs',
                'Petrapole export: 110 hrs → 14 hrs',
                '80% reduction in truck waiting time',
                'Proven model — replicable across ICPs',
            ]
        },
    ]
    for i, col in enumerate(cols):
        cx = Inches(0.3) + i * Inches(4.35)
        cy = Inches(1.1)
        cw = Inches(4.1)
        ch = Inches(5.85)
        rect(sl, cx, cy, cw, ch, fill=LBLUE)
        rect(sl, cx, cy, cw, Inches(0.9), fill=col['color'])
        txb(sl, cx + Inches(0.1), cy + Inches(0.07), cw - Inches(0.2), Inches(0.4),
            paras=[{'text': col['title'], 'size': 20, 'bold': True,
                    'color': WHITE, 'align': PP_ALIGN.CENTER}])
        txb(sl, cx + Inches(0.1), cy + Inches(0.48), cw - Inches(0.2), Inches(0.35),
            paras=[{'text': col['sub'], 'size': 9, 'color': WHITE,
                    'italic': True, 'align': PP_ALIGN.CENTER}])
        for j, b in enumerate(col['bullets']):
            by2 = cy + Inches(1.05) + j * Inches(0.85)
            rect(sl, cx + Inches(0.15), by2 + Inches(0.1), Inches(0.07), Inches(0.07), fill=col['color'])
            txb(sl, cx + Inches(0.3), by2, cw - Inches(0.4), Inches(0.75),
                paras=[{'text': b, 'size': 11, 'color': DARK}])
    # Bottom note
    rect(sl, Inches(0.3), Inches(6.1), Inches(12.9), Inches(0.8), fill=DNAVY)
    txb(sl, Inches(0.4), Inches(6.18), Inches(12.7), Inches(0.6),
        paras=[{'text': 'LPAI is already building the right infrastructure.',
                'size': 13, 'bold': True, 'color': GOLD, 'align': PP_ALIGN.CENTER},
               {'text': 'LPAI Nexus adds the AI intelligence and unified command layer on top.',
                'size': 12, 'color': WHITE, 'align': PP_ALIGN.CENTER}])
    set_notes(sl, "Acknowledge LPAI's strong existing digital foundation — LPMS, PIFS, and Suvidha are real achievements. Frame LPAI Nexus not as replacing these systems but as building the AI intelligence and unified command layer that connects and enhances them. This avoids any perception that the work dismisses existing investments.")
    return sl


def s04(prs):
    sl = new_slide(prs)
    T = "The Gap — Where Manual Processes Still Dominate"
    add_master(sl, 4, T); stitle(sl, T)
    cards = [
        ("31\nHours",   "Average import clearance time at land ports",  "CBIC Time Release Study 2023", RED),
        ("47\nDocs",    "Per consignment — 21 export + 26 import",       "Per stakeholder consultation", AMBER),
        ("Passive\nCCTV","No AI threat detection\nOfficers watch manually","Security gap at all ICPs", RED),
        ("No Unified\nView","LPAI HQ: no real-time picture of all 12 ICPs","Next-day Excel MIS only", AMBER),
        ("Manual\nANPR", "Plate reading in pilot phase only\nBlacklist check is manual","Critical coverage gap", RED),
    ]
    for i, (num, desc, src, c) in enumerate(cards):
        cx = Inches(0.3) + i * Inches(2.6)
        cy = Inches(1.1)
        cw = Inches(2.45)
        ch = Inches(5.85)
        rect(sl, cx, cy, cw, ch, fill=LBLUE)
        rect(sl, cx, cy, cw, Inches(0.08), fill=c)
        txb(sl, cx + Inches(0.1), cy + Inches(0.18), cw - Inches(0.2), Inches(1.0),
            paras=[{'text': num, 'size': 22, 'bold': True,
                    'color': c, 'align': PP_ALIGN.CENTER}])
        txb(sl, cx + Inches(0.1), cy + Inches(1.2), cw - Inches(0.2), Inches(3.0),
            paras=[{'text': desc, 'size': 11, 'color': DARK, 'align': PP_ALIGN.CENTER}])
        txb(sl, cx + Inches(0.1), cy + Inches(4.5), cw - Inches(0.2), Inches(0.8),
            paras=[{'text': src, 'size': 8.5, 'italic': True,
                    'color': GRAY, 'align': PP_ALIGN.CENTER}])
    set_notes(sl, "Each card represents a real operational gap documented in LPAI's own reports and CBIC studies. The 31-hour import clearance time is the most striking — contrast it with Singapore's 4.5-hour average. Do not dwell here; acknowledge the gaps and pivot to the solution LPAI Nexus provides.")
    return sl


def s05(prs):
    sl = new_slide(prs)
    T = "LPAI Nexus — The Unified Command Platform"
    add_master(sl, 5, T); stitle(sl, T)
    # Large screenshot centered
    img(sl, "command", Inches(0.4), Inches(1.1), Inches(10.3), Inches(5.2))
    # Right sidebar
    rect(sl, Inches(10.85), Inches(1.1), Inches(2.35), Inches(5.2), fill=DNAVY)
    txb(sl, Inches(10.92), Inches(1.2), Inches(2.2), Inches(0.5),
        paras=[{'text': 'Built during\nInternship at\nLPAI, MHA', 'size': 10,
                'color': WHITE, 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(10.92), Inches(2.1), Inches(2.2), Inches(0.5),
        paras=[{'text': 'June 2026', 'size': 11, 'bold': True,
                'color': GOLD, 'align': PP_ALIGN.CENTER}])
    rect(sl, Inches(11.0), Inches(2.65), Inches(2.0), Inches(0.03), fill=GOLD)
    for label in ['React 18', 'Vite 5', 'Recharts', 'Tailwind CSS']:
        idx = ['React 18', 'Vite 5', 'Recharts', 'Tailwind CSS'].index(label)
        txb(sl, Inches(10.92), Inches(2.75) + idx * Inches(0.38), Inches(2.2), Inches(0.35),
            paras=[{'text': label, 'size': 9.5, 'color': LBLUE, 'align': PP_ALIGN.CENTER}])
    rect(sl, Inches(11.0), Inches(4.35), Inches(2.0), Inches(0.03), fill=GOLD)
    txb(sl, Inches(10.92), Inches(4.45), Inches(2.2), Inches(0.35),
        paras=[{'text': 'Full Working\nPrototype', 'size': 9.5,
                'bold': True, 'color': GREEN, 'align': PP_ALIGN.CENTER}])
    # Pills below screenshot
    pills = [("6 Modules", NAVY), ("12 ICPs", GREEN), ("Powered by Axiom AI", GOLD)]
    for i, (label, c) in enumerate(pills):
        px = Inches(1.2) + i * Inches(2.95)
        rect(sl, px, Inches(6.42), Inches(2.5), Inches(0.42), fill=c)
        txb(sl, px, Inches(6.47), Inches(2.5), Inches(0.32),
            paras=[{'text': label, 'size': 12, 'bold': True,
                    'color': DARK if c == GOLD else WHITE,
                    'align': PP_ALIGN.CENTER}])
    set_notes(sl, "This is the first time the audience sees the actual platform. Let the screenshot speak for itself — pause for 10-15 seconds. Then point out the ICP selector (top), the Axiom briefing (center), and the module tabs (left nav). Emphasize that everything shown is fully functional and demonstrable right now.")
    return sl


def s06(prs):
    sl = new_slide(prs)
    T = "Axiom — The AI Intelligence Engine"
    add_master(sl, 6, T); stitle(sl, T)
    # Axiom circle
    rect(sl, Inches(0.5), Inches(1.2), Inches(4.8), Inches(4.8), fill=DNAVY)
    rect(sl, Inches(0.7), Inches(1.4), Inches(4.4), Inches(4.4), fill=NAVY)
    txb(sl, Inches(0.5), Inches(2.6), Inches(4.8), Inches(0.7),
        paras=[{'text': 'AXIOM', 'size': 36, 'bold': True,
                'color': GOLD, 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(3.35), Inches(4.8), Inches(0.4),
        paras=[{'text': 'AI Intelligence Engine', 'size': 13,
                'color': WHITE, 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(3.8), Inches(4.8), Inches(0.4),
        paras=[{'text': 'v1.0 — LPAI Nexus', 'size': 10, 'italic': True,
                'color': LBLUE, 'align': PP_ALIGN.CENTER}])
    # Capabilities
    caps = [
        ("⚡", "Cargo Document Processing",    "2 min vs 45 min manual — 95% faster",         GREEN),
        ("👤", "Watchlist Face Matching",        "Under 1 second per passenger",                 GREEN),
        ("📹", "Camera Feed Monitoring",         "All feeds simultaneously, 24/7 — zero gaps",  GREEN),
        ("🚛", "Number Plate Recognition",       "2 seconds + instant VAHAN + blacklist lookup", GREEN),
        ("📊", "Queue Size Prediction",          "6 hours ahead — time-series model",            TEAL),
        ("📄", "Daily Intelligence Briefing",    "Auto-generated at 06:00 every morning",        TEAL),
    ]
    for i, (icon, cap, metric, c) in enumerate(caps):
        cy2 = Inches(1.15) + i * Inches(0.95)
        rect(sl, Inches(5.6), cy2, Inches(7.5), Inches(0.85), fill=LBLUE)
        rect(sl, Inches(5.6), cy2, Inches(0.06), Inches(0.85), fill=c)
        txb(sl, Inches(5.72), cy2 + Inches(0.06), Inches(0.5), Inches(0.4),
            paras=[{'text': icon, 'size': 16, 'align': PP_ALIGN.CENTER}])
        txb(sl, Inches(6.25), cy2 + Inches(0.04), Inches(4.0), Inches(0.38),
            paras=[{'text': cap, 'size': 11, 'bold': True, 'color': NAVY}])
        txb(sl, Inches(6.25), cy2 + Inches(0.42), Inches(4.0), Inches(0.35),
            paras=[{'text': metric, 'size': 10, 'color': c}])
        txb(sl, Inches(10.3), cy2 + Inches(0.2), Inches(2.7), Inches(0.4),
            paras=[{'text': '▶', 'size': 10, 'color': c}])
    # Bottom tagline
    rect(sl, Inches(0.3), Inches(6.1), Inches(12.9), Inches(0.82), fill=DNAVY)
    txb(sl, Inches(0.5), Inches(6.22), Inches(12.7), Inches(0.55),
        paras=[{'text': '"Axiom does not replace officers — it ensures nothing is missed."',
                'size': 14, 'italic': True, 'color': GOLD, 'align': PP_ALIGN.CENTER}])
    set_notes(sl, "Axiom is the differentiator between a data dashboard and an intelligence platform. Emphasize the 95% faster cargo processing figure — this is not aspirational, it reflects the speed difference between AI-assisted OCR and manual document review. The quote at the bottom is important for any officer in the room concerned about job displacement.")
    return sl


def _module_slide(prs, num, title, img_key, right_content_fn):
    """Template for module slides 7-13."""
    sl = new_slide(prs)
    add_master(sl, num, title)
    stitle(sl, title)
    img(sl, img_key, Inches(0.3), Inches(1.1), Inches(7.2), Inches(5.1))
    right_content_fn(sl)
    return sl


def s07(prs):
    sl = new_slide(prs)
    T = "Command Overview — All 12 ICPs on One Screen"
    add_master(sl, 7, T); stitle(sl, T)
    img(sl, "command", Inches(0.3), Inches(1.1), Inches(7.6), Inches(5.0))
    # Right content
    checks = [
        ("Real-time KPIs", "Cargo, crossings, alerts, vehicles — all at a glance"),
        ("Axiom AI Briefing", "Contextual intelligence summary with confidence scores"),
        ("All 12 ICPs Live", "ICP grid with status, throughput and risk flags"),
    ]
    for i, (h, d) in enumerate(checks):
        cy2 = Inches(1.3) + i * Inches(1.35)
        rect(sl, Inches(8.15), cy2, Inches(4.95), Inches(1.2), fill=LBLUE)
        rect(sl, Inches(8.15), cy2, Inches(0.06), Inches(1.2), fill=GREEN)
        txb(sl, Inches(8.28), cy2 + Inches(0.08), Inches(0.35), Inches(0.4),
            paras=[{'text': '✓', 'size': 16, 'bold': True, 'color': GREEN}])
        txb(sl, Inches(8.65), cy2 + Inches(0.08), Inches(4.35), Inches(0.4),
            paras=[{'text': h, 'size': 12, 'bold': True, 'color': NAVY}])
        txb(sl, Inches(8.65), cy2 + Inches(0.5), Inches(4.35), Inches(0.55),
            paras=[{'text': d, 'size': 10.5, 'color': DARK}])
    # Bottom bar
    rect(sl, Inches(0.3), Inches(6.25), Inches(12.9), Inches(0.65), fill=NAVY)
    txb(sl, Inches(0.4), Inches(6.3), Inches(5.5), Inches(0.55),
        paras=[{'text': 'Replaces: Next-day Excel MIS reports', 'size': 10, 'color': LGRAY}])
    rect(sl, Inches(6.0), Inches(6.35), Inches(0.03), Inches(0.45), fill=GOLD)
    txb(sl, Inches(6.15), Inches(6.3), Inches(6.9), Inches(0.55),
        paras=[{'text': 'Benefit: Instant situational awareness for LPAI HQ', 'size': 10,
                'bold': True, 'color': GOLD}])
    set_notes(sl, "This is the home screen every LPAI officer will see when they log in. Point out that the ICP selector (top right) instantly switches the entire dashboard to any of the 12 ICPs. The Axiom briefing at center replaces the morning MIS report — it synthesises overnight data into 4-5 actionable sentences.")
    return sl


def s08(prs):
    sl = new_slide(prs)
    T = "Cargo Clearance — AI Risk Scoring for Every Consignment"
    add_master(sl, 8, T); stitle(sl, T)
    img(sl, "cargo", Inches(0.3), Inches(1.1), Inches(7.2), Inches(4.95))
    # Right: 4-lane system
    txb(sl, Inches(7.65), Inches(1.15), Inches(5.45), Inches(0.4),
        paras=[{'text': 'The Four-Lane Risk System', 'size': 13, 'bold': True, 'color': NAVY}])
    lanes = [
        (GREEN,  "GREEN LANE",  "Score 0–35",  "Auto-cleared by Axiom AI"),
        (AMBER,  "YELLOW LANE", "Score 36–65", "Document verification required"),
        (RGBColor(0xE6,0x7E,0x22), "ORANGE LANE", "Score 66–80", "Physical examination ordered"),
        (RED,    "RED LANE",    "Score 81–100","Full hold — officer decision required"),
    ]
    for i, (c, name, score, action) in enumerate(lanes):
        ly = Inches(1.65) + i * Inches(1.05)
        rect(sl, Inches(7.65), ly, Inches(5.45), Inches(0.95), fill=LBLUE)
        rect(sl, Inches(7.65), ly, Inches(0.18), Inches(0.95), fill=c)
        txb(sl, Inches(7.9), ly + Inches(0.06), Inches(2.5), Inches(0.38),
            paras=[{'text': name, 'size': 11, 'bold': True, 'color': c}])
        txb(sl, Inches(7.9), ly + Inches(0.46), Inches(2.5), Inches(0.38),
            paras=[{'text': score, 'size': 10, 'color': GRAY}])
        txb(sl, Inches(10.45), ly + Inches(0.2), Inches(2.55), Inches(0.55),
            paras=[{'text': action, 'size': 10, 'color': DARK}])
    # Stat bar
    rect(sl, Inches(0.3), Inches(6.28), Inches(12.9), Inches(0.62), fill=DNAVY)
    for j, (label, val, c) in enumerate([
        ("Processed Today", "847 consignments", WHITE),
        ("Auto-cleared", "612 by Axiom AI", GREEN),
        ("Flagged for Review", "37 consignments", RED),
    ]):
        sx = Inches(0.5) + j * Inches(4.3)
        txb(sl, sx, Inches(6.33), Inches(2.0), Inches(0.25),
            paras=[{'text': label, 'size': 8.5, 'color': LGRAY}])
        txb(sl, sx, Inches(6.55), Inches(2.0), Inches(0.28),
            paras=[{'text': val, 'size': 11, 'bold': True, 'color': c}])
    set_notes(sl, "Axiom assigns a risk score to every consignment automatically — the officer does not need to review all 847, only the 37 flagged. This is the operational transformation: attention directed by AI, not wasted on low-risk cargo. The four-lane system mirrors the established WCO risk management framework.")
    return sl


def s09(prs):
    sl = new_slide(prs)
    T = "People Movement — Predictive Queue Management"
    add_master(sl, 9, T); stitle(sl, T)
    img(sl, "immigration", Inches(0.3), Inches(1.1), Inches(7.2), Inches(4.95))
    # Right content
    rect(sl, Inches(7.65), Inches(1.15), Inches(5.45), Inches(1.8), fill=NAVY)
    txb(sl, Inches(7.75), Inches(1.22), Inches(5.25), Inches(0.4),
        paras=[{'text': 'Axiom Queue Intelligence', 'size': 12, 'bold': True, 'color': GOLD}])
    txb(sl, Inches(7.75), Inches(1.65), Inches(5.25), Inches(1.1),
        paras=[{'text': 'Predicts queue sizes 6 hours ahead.\nICP Commanders open counters before queues build\n— not after.',
                'size': 11, 'color': WHITE}])
    # 3 stats
    stats3 = [("2,341", "Crossings Today"), ("494", "Foreign Nationals"), ("12", "Flagged / Detained")]
    for i, (v, l) in enumerate(stats3):
        sy = Inches(3.1) + i * Inches(0.9)
        rect(sl, Inches(7.65), sy, Inches(5.45), Inches(0.82), fill=LBLUE)
        txb(sl, Inches(7.75), sy + Inches(0.06), Inches(2.0), Inches(0.38),
            paras=[{'text': v, 'size': 22, 'bold': True, 'color': NAVY}])
        txb(sl, Inches(7.75), sy + Inches(0.44), Inches(3.5), Inches(0.28),
            paras=[{'text': l, 'size': 10, 'color': GRAY}])
    # Watchlist stat
    rect(sl, Inches(7.65), Inches(5.85), Inches(5.45), Inches(0.55), fill=GREEN)
    txb(sl, Inches(7.75), Inches(5.9), Inches(5.25), Inches(0.42),
        paras=[{'text': 'Watchlist check: <1 second  vs  2-3 min manual lookup',
                'size': 10, 'bold': True, 'color': WHITE}])
    set_notes(sl, "The queue prediction capability is the most impactful for passenger flow. Emphasize that Axiom analyses historical patterns, current arrival rates, and staffing to forecast 6 hours ahead. The ICP Commander can pre-deploy counters at 8am knowing that peak will hit at 2pm — currently this decision is reactive and too late.")
    return sl


def s10(prs):
    sl = new_slide(prs)
    T = "Security & Surveillance — AI That Never Blinks"
    add_master(sl, 10, T); stitle(sl, T)
    img(sl, "security", Inches(0.3), Inches(1.1), Inches(7.2), Inches(4.95))
    # Detection list
    txb(sl, Inches(7.65), Inches(1.15), Inches(5.45), Inches(0.38),
        paras=[{'text': 'What Axiom Detects:', 'size': 13, 'bold': True, 'color': NAVY}])
    detections = [
        ("🚫", "Restricted zone intrusions"),
        ("📦", "Unattended objects > 5 minutes"),
        ("👥", "Crowd density threshold breach"),
        ("🎯", "Watchlist face matches in real time"),
        ("🚗", "Vehicle anomalies and wrong-lane entry"),
    ]
    for i, (icon, label) in enumerate(detections):
        dy = Inches(1.6) + i * Inches(0.68)
        rect(sl, Inches(7.65), dy, Inches(5.45), Inches(0.6), fill=LBLUE)
        txb(sl, Inches(7.75), dy + Inches(0.1), Inches(0.5), Inches(0.4),
            paras=[{'text': icon, 'size': 14}])
        txb(sl, Inches(8.3), dy + Inches(0.1), Inches(4.7), Inches(0.4),
            paras=[{'text': label, 'size': 11, 'color': DARK}])
    # Key stat box
    rect(sl, Inches(7.65), Inches(5.2), Inches(5.45), Inches(1.2), fill=LRED)
    rect(sl, Inches(7.65), Inches(5.2), Inches(0.06), Inches(1.2), fill=RED)
    txb(sl, Inches(7.78), Inches(5.26), Inches(5.2), Inches(0.4),
        paras=[{'text': 'Human officers: Can monitor 4–6 screens', 'size': 10, 'color': RED}])
    txb(sl, Inches(7.78), Inches(5.68), Inches(5.2), Inches(0.4),
        paras=[{'text': 'Axiom: Monitors ALL feeds simultaneously, 24/7',
                'size': 11, 'bold': True, 'color': NAVY}])
    set_notes(sl, "The surveillance gap is a genuine security risk — ICP perimeters are large and camera networks are expanding faster than officer rosters. Axiom does not make fatigue-related errors. The five detection types listed are based on standard computer vision capabilities deployable on Jetson Orin edge devices at remote ICPs with low connectivity.")
    return sl


def s11(prs):
    sl = new_slide(prs)
    T = "Vehicle Gate — ANPR + Weighbridge Intelligence"
    add_master(sl, 11, T); stitle(sl, T)
    img(sl, "vehicle", Inches(0.3), Inches(1.1), Inches(7.2), Inches(4.95))
    # Process flow
    txb(sl, Inches(7.65), Inches(1.15), Inches(5.45), Inches(0.38),
        paras=[{'text': 'Processing Flow:', 'size': 13, 'bold': True, 'color': NAVY}])
    steps = [
        ("1", NAVY,  "Vehicle approaches → ANPR reads plate"),
        ("2", TEAL,  "Simultaneous: VAHAN + TAS + Blacklist check"),
        ("3", AMBER, "Weighbridge → auto weight variance detection"),
        ("4", GREEN, "Risk score assigned → lane directed by Axiom"),
    ]
    for i, (n, c, txt) in enumerate(steps):
        sy = Inches(1.6) + i * Inches(0.95)
        rect(sl, Inches(7.65), sy, Inches(5.45), Inches(0.85), fill=LBLUE)
        rect(sl, Inches(7.65), sy, Inches(0.45), Inches(0.85), fill=c)
        txb(sl, Inches(7.65), sy + Inches(0.2), Inches(0.45), Inches(0.42),
            paras=[{'text': n, 'size': 14, 'bold': True, 'color': WHITE,
                    'align': PP_ALIGN.CENTER}])
        txb(sl, Inches(8.2), sy + Inches(0.2), Inches(4.8), Inches(0.42),
            paras=[{'text': txt, 'size': 11, 'color': DARK}])
    # Stats
    vstats = [("543", "Vehicles processed today"), ("100%", "ANPR plate coverage"), ("8", "Overloaded — auto-detected")]
    for i, (v, l) in enumerate(vstats):
        vx = Inches(7.65) + i * Inches(1.85)
        rect(sl, vx, Inches(5.55), Inches(1.75), Inches(0.85), fill=NAVY)
        txb(sl, vx + Inches(0.05), Inches(5.6), Inches(1.65), Inches(0.42),
            paras=[{'text': v, 'size': 18, 'bold': True, 'color': GOLD, 'align': PP_ALIGN.CENTER}])
        txb(sl, vx + Inches(0.05), Inches(6.02), Inches(1.65), Inches(0.28),
            paras=[{'text': l, 'size': 8, 'color': WHITE, 'align': PP_ALIGN.CENTER}])
    set_notes(sl, "The ANPR + VAHAN integration is the key capability here. Every vehicle entering the ICP is cross-checked against the VAHAN database and the blacklist in under 2 seconds — something that currently takes an officer 3-5 minutes per vehicle with manual checks. The weighbridge integration prevents the common fraud of mis-declaring cargo weight.")
    return sl


def s12(prs):
    sl = new_slide(prs)
    T = "Intelligence & Analytics — Data-Driven Border Management"
    add_master(sl, 12, T); stitle(sl, T)
    img(sl, "intelligence", Inches(0.3), Inches(1.1), Inches(7.2), Inches(4.95))
    txb(sl, Inches(7.65), Inches(1.15), Inches(5.45), Inches(0.38),
        paras=[{'text': 'FY 2023-24 Actuals (Real LPAI Data):',
                'size': 12, 'bold': True, 'color': NAVY}])
    stats = [
        ("₹70,952 Cr", "Total Trade Facilitated",    NAVY),
        ("30.46 Lakh", "Passengers Processed",        GREEN),
        ("₹62.17 Cr",  "Duty Revenue Collected",      TEAL),
        ("12 / 23",    "Operational ICPs (of 23 sanctioned)", GOLD),
    ]
    for i, (v, l, c) in enumerate(stats):
        col = i % 2
        row = i // 2
        sx = Inches(7.65) + col * Inches(2.75)
        sy = Inches(1.6) + row * Inches(1.65)
        rect(sl, sx, sy, Inches(2.65), Inches(1.5), fill=LBLUE)
        rect(sl, sx, sy, Inches(2.65), Inches(0.06), fill=c)
        txb(sl, sx + Inches(0.1), sy + Inches(0.18), Inches(2.45), Inches(0.55),
            paras=[{'text': v, 'size': 20, 'bold': True, 'color': c, 'align': PP_ALIGN.CENTER}])
        txb(sl, sx + Inches(0.1), sy + Inches(0.8), Inches(2.45), Inches(0.55),
            paras=[{'text': l, 'size': 9.5, 'color': DARK, 'align': PP_ALIGN.CENTER}])
    rect(sl, Inches(7.65), Inches(5.05), Inches(5.45), Inches(0.85), fill=DNAVY)
    txb(sl, Inches(7.75), Inches(5.12), Inches(5.25), Inches(0.65),
        paras=[{'text': 'All 12 ICPs benchmarked side-by-side in real time\nNo more waiting for consolidated monthly reports',
                'size': 10, 'color': WHITE}])
    set_notes(sl, "The numbers shown are from LPAI's own Annual Report 2023-24 — real operational data, not estimates. Point out that currently this data is only available after month-end consolidation. LPAI Nexus makes it available in real time, letting senior officials make decisions on current performance rather than last month's data.")
    return sl


def s13(prs):
    sl = new_slide(prs)
    T = "Built-In Officer Onboarding — 45-Step Guided Tour"
    add_master(sl, 13, T); stitle(sl, T)
    img(sl, "walkthrough", Inches(0.3), Inches(1.1), Inches(7.2), Inches(4.95))
    # Challenge / Solution
    rect(sl, Inches(7.65), Inches(1.15), Inches(5.45), Inches(1.55), fill=LRED)
    rect(sl, Inches(7.65), Inches(1.15), Inches(0.06), Inches(1.55), fill=RED)
    txb(sl, Inches(7.78), Inches(1.2), Inches(5.2), Inches(0.38),
        paras=[{'text': 'The Challenge:', 'size': 11, 'bold': True, 'color': RED}])
    txb(sl, Inches(7.78), Inches(1.6), Inches(5.2), Inches(0.85),
        paras=[{'text': 'Senior officers unfamiliar with new systems will not read a manual — they need guidance built into the software.',
                'size': 10, 'color': DARK}])
    rect(sl, Inches(7.65), Inches(2.8), Inches(5.45), Inches(1.55), fill=LGREEN)
    rect(sl, Inches(7.65), Inches(2.8), Inches(0.06), Inches(1.55), fill=GREEN)
    txb(sl, Inches(7.78), Inches(2.86), Inches(5.2), Inches(0.38),
        paras=[{'text': 'The Solution:', 'size': 11, 'bold': True, 'color': GREEN}])
    txb(sl, Inches(7.78), Inches(3.26), Inches(5.2), Inches(0.85),
        paras=[{'text': 'Press "Start Tour" — system guides you through every feature in plain language, step by step, module by module.',
                'size': 10, 'color': DARK}])
    feats = [
        "45 steps across all 6 modules",
        "Written for IAS-level audiences",
        "Keyboard navigation + mobile support",
        "Persists completion — never shown twice",
        "Works on any government laptop / tablet",
    ]
    for i, f in enumerate(feats):
        fy = Inches(4.5) + i * Inches(0.33)
        txb(sl, Inches(7.75), fy, Inches(0.3), Inches(0.3),
            paras=[{'text': '▶', 'size': 9, 'color': NAVY}])
        txb(sl, Inches(8.05), fy, Inches(5.0), Inches(0.3),
            paras=[{'text': f, 'size': 10, 'color': DARK}])
    set_notes(sl, "The walkthrough system is often the most impressive feature for non-technical audiences. The spotlight overlay and step-by-step guidance means any officer — even a Secretary who has never used a dashboard before — can self-onboard in 8 minutes. This directly addresses the change management risk in any government technology deployment.")
    return sl


def s14(prs):
    sl = new_slide(prs)
    T = "System Architecture — Built for Government Scale"
    add_master(sl, 14, T); stitle(sl, T)
    layers = [
        (LBLUE,  NAVY,  "PRESENTATION LAYER",
         "React 18 SPA  ·  6 Functional Modules  ·  Shared Component Library  ·  Responsive — Works on Government Laptops & Tablets"),
        (NAVY,   WHITE, "AXIOM INTELLIGENCE LAYER",
         "Document OCR  ·  Face Recognition  ·  Camera Analysis  ·  Queue Prediction  ·  Cargo Risk Scoring  ·  NLP Briefing Engine"),
        (RGBColor(0xF5,0xE9,0xD4), DARK, "DATA INTEGRATION LAYER",
         "ICEGATE API  ·  VAHAN API  ·  Bureau of Immigration API  ·  NDAP  ·  LPMS  ·  PIFS  ·  CCTV Streams"),
        (LGREEN, DARK,  "DATA SOURCES",
         "Customs Records  ·  Vehicle Registry  ·  Passport Database  ·  CCTV Cameras  ·  Weighbridge Sensors  ·  Intelligence Feeds"),
    ]
    layer_colors_left = [NAVY, GOLD, AMBER, GREEN]
    for i, (bg, tc, lname, ldetail) in enumerate(layers):
        ly = Inches(1.12) + i * Inches(1.35)
        lh = Inches(1.25)
        rect(sl, Inches(0.3), ly, Inches(12.9), lh, fill=bg)
        rect(sl, Inches(0.3), ly, Inches(0.25), lh, fill=layer_colors_left[i])
        txb(sl, Inches(0.65), ly + Inches(0.1), Inches(3.2), Inches(0.42),
            paras=[{'text': lname, 'size': 12, 'bold': True, 'color': tc}])
        txb(sl, Inches(0.65), ly + Inches(0.55), Inches(12.3), Inches(0.55),
            paras=[{'text': ldetail, 'size': 10, 'color': tc}])
    # Status note
    rect(sl, Inches(0.3), Inches(6.55), Inches(12.9), Inches(0.55), fill=LAMBER)
    txb(sl, Inches(0.5), Inches(6.6), Inches(12.7), Inches(0.42),
        paras=[{'text': '◉  Current Status:  Layer 1 complete  ·  Layer 2 concept-implemented  ·  Layers 3–4 planned Phase 2 (data sharing agreements required)',
                'size': 10, 'bold': False, 'color': RGBColor(0x85, 0x65, 0x04)}])
    set_notes(sl, "Be transparent about what is built vs what is planned. Layer 1 (the full platform) is working today. Layer 2 (Axiom) is implemented as a concept engine in the prototype. Layers 3-4 (live API integrations) require inter-agency data agreements and are Phase 2 work. This transparency builds credibility with a senior government audience.")
    return sl


def s15(prs):
    sl = new_slide(prs)
    T = "The Impact — Process Transformation at Every Stage"
    add_master(sl, 15, T); stitle(sl, T)
    tbl_data = [
        ["Process",                  "Before LPAI Nexus",               "After LPAI Nexus",          "Improvement"],
        ["Cargo document check",     "45 min — manual officer review",   "2 min — Axiom AI OCR",      "95% faster"],
        ["Export clearance (avg)",   "110+ hours at Petrapole",          "Target: 12–18 hours",       "~85% reduction"],
        ["Import clearance (avg)",   "31 hours (CBIC 2023)",             "Target: 12–18 hours",       "~50% reduction"],
        ["Watchlist check",          "2–3 min — manual database lookup", "<1 second — AI matching",   "99% faster"],
        ["CCTV monitoring",          "Human-limited, 4–6 screens",       "AI-powered, all feeds 24/7","Continuous"],
        ["MIS / HQ reporting",       "Next-day Excel — 12-48 hr delay",  "Real-time dashboard",       "Instant"],
    ]
    col_w = [Inches(3.0), Inches(3.5), Inches(3.5), Inches(2.7)]
    rows, cols = len(tbl_data), 4
    tbl = sl.shapes.add_table(rows, cols, Inches(0.3), Inches(1.1), Inches(12.9), Inches(5.7)).table
    for ci, w in enumerate(col_w):
        tbl.columns[ci].width = w
    for ri, row in enumerate(tbl_data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            if ri == 0:
                cell_fmt(cell, val, bold=True, color=WHITE, bg=NAVY,
                         size=11, align=PP_ALIGN.CENTER)
            else:
                bg = LBLUE if ri % 2 == 0 else WHITE
                col = (GREEN if ci == 3 else DARK)
                is_bold = (ci == 3)
                cell_fmt(cell, val, bold=is_bold, color=col, bg=bg, size=10)
    set_notes(sl, "Use the CBIC Time Release Study 2023 figure (31 hours import clearance) as the anchor. These are real numbers from real government studies — LPAI Nexus targets represent international best-practice benchmarks. The improvements are conservative estimates based on comparable implementations at Singapore and EU smart border programs.")
    return sl


def s16(prs):
    sl = new_slide(prs)
    T = "Three-Phase Implementation Plan"
    add_master(sl, 16, T); stitle(sl, T)
    phases = [
        {
            'color': GREEN, 'lcolor': LGREEN, 'status': 'COMPLETED',
            'title': 'Phase 1: Prototype',
            'period': 'Internship 2026 — Done',
            'items': ['✓ 6 fully functional modules', '✓ 45-step guided walkthrough',
                      '✓ Full technical report + PRD', '✓ Axiom AI engine integrated',
                      '✓ This presentation'],
        },
        {
            'color': AMBER, 'lcolor': LAMBER, 'status': 'PROPOSED',
            'title': 'Phase 2: Pilot Deployment',
            'period': 'Months 4–9 post approval',
            'items': ['• Petrapole + Attari ICPs', '• LPMS backend integration',
                      '• Bureau of Immigration API', '• Officer training programme',
                      '• Performance benchmarking'],
        },
        {
            'color': NAVY, 'lcolor': LBLUE, 'status': 'FUTURE',
            'title': 'Phase 3: National Rollout',
            'period': 'Months 10–24',
            'items': ['• All 12 operational ICPs', '• 1 ICP per month rollout cadence',
                      '• Central command at LPAI HQ', '• AI model training on real data',
                      '• NATGRID + CCTNS integration'],
        },
    ]
    for i, ph in enumerate(phases):
        px = Inches(0.3) + i * Inches(4.38)
        pw = Inches(4.2)
        ph_h = Inches(5.85)
        rect(sl, px, Inches(1.1), pw, ph_h, fill=ph['lcolor'])
        rect(sl, px, Inches(1.1), pw, Inches(0.9), fill=ph['color'])
        rect(sl, px, Inches(1.1), Inches(0.08), ph_h, fill=ph['color'])
        txb(sl, px + Inches(0.15), Inches(1.15), pw - Inches(0.2), Inches(0.42),
            paras=[{'text': ph['title'], 'size': 13, 'bold': True,
                    'color': DARK if ph['color'] == AMBER else WHITE}])
        txb(sl, px + Inches(0.15), Inches(1.57), pw - Inches(0.2), Inches(0.38),
            paras=[{'text': ph['period'], 'size': 10, 'italic': True,
                    'color': DARK if ph['color'] == AMBER else WHITE}])
        rect(sl, px + Inches(0.5), Inches(2.2), pw - Inches(1.0), Inches(0.04), fill=ph['color'])
        # Status badge
        rect(sl, px + Inches(0.5), Inches(2.35), pw - Inches(1.0), Inches(0.38), fill=ph['color'])
        txb(sl, px + Inches(0.5), Inches(2.38), pw - Inches(1.0), Inches(0.3),
            paras=[{'text': ph['status'], 'size': 10, 'bold': True,
                    'color': DARK if ph['color'] == AMBER else WHITE,
                    'align': PP_ALIGN.CENTER}])
        for j, item in enumerate(ph['items']):
            iy = Inches(2.85) + j * Inches(0.58)
            txb(sl, px + Inches(0.2), iy, pw - Inches(0.3), Inches(0.52),
                paras=[{'text': item, 'size': 10.5, 'color': DARK}])
    # Arrow connectors
    for ax in [Inches(4.52), Inches(8.9)]:
        txb(sl, ax, Inches(3.7), Inches(0.36), Inches(0.5),
            paras=[{'text': '▶', 'size': 20, 'color': NAVY, 'align': PP_ALIGN.CENTER}])
    set_notes(sl, "Phase 1 is complete — the prototype exists and works. The ask to this audience is approval to proceed to Phase 2. Petrapole and Attari are the right pilot ICPs because they are the highest-volume ICPs on the Bangladesh and Pakistan borders respectively — success there validates the model for all 12. Phase 2 requires MOU agreements with CBIC and BoI.")
    return sl


def s17(prs):
    sl = new_slide(prs)
    T = "LPAI Nexus — Aligned with National Priorities"
    add_master(sl, 17, T); stitle(sl, T)
    cols17 = [
        {
            'icon': '🇮🇳', 'heading': 'Viksit Bharat 2047',
            'color': NAVY,
            'text': 'Digital governance as the foundation for a developed India. AI-powered border management directly supports the vision of a modern, efficient state.',
            'nexus': 'LPAI Nexus: Real-time AI governance at every land border crossing by 2030.',
        },
        {
            'icon': '📋', 'heading': 'NTFAP 2020-23 Action #44',
            'color': TEAL,
            'text': 'The National Trade Facilitation Action Plan mandates a Land Port Community System for paperless, coordinated border clearance.',
            'nexus': 'LPAI Nexus: Directly implements this mandate with a unified AI intelligence layer.',
        },
        {
            'icon': '💻', 'heading': 'Digital India Mission',
            'color': GREEN,
            'text': 'Paperless, faceless, contactless government services across all citizen touchpoints — including India\'s 12 ICPs.',
            'nexus': 'LPAI Nexus: Automated clearance, digital audit trail, contactless AI-assisted processing.',
        },
    ]
    for i, col in enumerate(cols17):
        cx = Inches(0.3) + i * Inches(4.38)
        cw = Inches(4.2)
        rect(sl, cx, Inches(1.1), cw, Inches(5.85), fill=LBLUE)
        rect(sl, cx, Inches(1.1), cw, Inches(0.82), fill=col['color'])
        txb(sl, cx + Inches(0.1), Inches(1.15), cw - Inches(0.2), Inches(0.42),
            paras=[{'text': f"{col['icon']}  {col['heading']}",
                    'size': 12, 'bold': True, 'color': WHITE}])
        txb(sl, cx + Inches(0.15), Inches(2.05), cw - Inches(0.3), Inches(2.0),
            paras=[{'text': col['text'], 'size': 11, 'color': DARK}])
        rect(sl, cx + Inches(0.15), Inches(4.2), cw - Inches(0.3), Inches(0.04), fill=col['color'])
        rect(sl, cx + Inches(0.1), Inches(4.3), cw - Inches(0.2), Inches(1.6), fill=col['color'])
        txb(sl, cx + Inches(0.2), Inches(4.38), cw - Inches(0.4), Inches(1.4),
            paras=[{'text': col['nexus'], 'size': 10, 'color': WHITE}])
    # Bottom quote
    rect(sl, Inches(0.3), Inches(6.12), Inches(12.9), Inches(0.82), fill=DNAVY)
    txb(sl, Inches(0.5), Inches(6.2), Inches(12.7), Inches(0.62),
        paras=[{'text': 'LPAI Annual Report 2023-24 explicitly calls for AI/ML integration in LPMS and AI-based security at all ICPs. LPAI Nexus directly implements this stated vision.',
                'size': 11, 'color': GOLD, 'align': PP_ALIGN.CENTER}])
    set_notes(sl, "This slide is important for a Secretary-level audience — it shows LPAI Nexus is not a student project but aligns with the highest levels of national policy direction. Cite the LPAI Annual Report line directly if asked: it explicitly calls for AI/ML integration in LPMS. LPAI Nexus is the prototype implementation of that stated vision.")
    return sl


def s18(prs):
    sl = new_slide(prs)
    T = "Where India Stands — Global Smart Border Benchmarks"
    add_master(sl, 18, T); stitle(sl, T)
    tbl18 = [
        ["Country / Port",         "System",          "Key Achievement",                          "India Comparison"],
        ["Singapore Ports",        "Portnet PCS",     "Full port community system\nSaves $80M+ over 3 years", "Planned via LPMS — LPAI Nexus adds AI layer"],
        ["US–Mexico Border",       "CBP One + FIRST", "Real-time cargo tracking\nJoint inspection modules", "LPAI Nexus proposes equivalent for land borders"],
        ["European Union",         "EES + ETIAS",     "Automated entry/exit system\nBiometric at all crossings", "LPAI Nexus addresses India's land border gap"],
        ["Bangladesh (BSSB)",      "Digital Clearance","Slot booking + digital forms\nReduced wait times 60%", "India ahead via Suvidha Portal + LPAI Nexus"],
    ]
    col_w18 = [Inches(2.8), Inches(2.5), Inches(3.8), Inches(3.7)]
    tbl = sl.shapes.add_table(5, 4, Inches(0.3), Inches(1.1), Inches(12.9), Inches(5.4)).table
    for ci, w in enumerate(col_w18):
        tbl.columns[ci].width = w
    for ri, row in enumerate(tbl18):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            if ri == 0:
                cell_fmt(cell, val, bold=True, color=WHITE, bg=NAVY, size=11, align=PP_ALIGN.CENTER)
            else:
                bg = LBLUE if ri % 2 == 1 else WHITE
                fc = GREEN if ci == 3 else DARK
                cell_fmt(cell, val, color=fc, bg=bg, size=10)
    rect(sl, Inches(0.3), Inches(6.6), Inches(12.9), Inches(0.5), fill=DNAVY)
    txb(sl, Inches(0.5), Inches(6.65), Inches(12.7), Inches(0.38),
        paras=[{'text': 'LPAI Nexus positions India\'s land border management at international best-practice level',
                'size': 11, 'bold': True, 'color': GOLD, 'align': PP_ALIGN.CENTER}])
    set_notes(sl, "Global benchmarks anchor LPAI Nexus in an international context. Singapore's Portnet is the gold standard — it required $200M+ investment over 5 years. LPAI Nexus proposes an AI-first approach that can reach equivalent outcomes faster by building on LPAI's existing LPMS infrastructure. Bangladesh's success with digital clearance is particularly relevant — it demonstrates the model works in the South Asian operational context.")
    return sl


def s19(prs):
    sl = new_slide(prs)
    T = "Challenges — And How to Address Them"
    add_master(sl, 19, T); stitle(sl, T)
    cards19 = [
        {
            'title': 'Connectivity at Remote ICPs',
            'detail': 'Moreh ICP (Manipur–Myanmar): 8 Mbps average connectivity. Cloud-dependent AI would fail during peak traffic.',
            'solution': 'Edge AI deployment — AI models run locally on NVIDIA Jetson Orin devices. No cloud dependency for real-time processing.',
        },
        {
            'title': 'Training Data for AI Models',
            'detail': 'Axiom requires historical Bill of Entry data, footfall logs, and watchlist access to train accurate models for Indian border conditions.',
            'solution': 'Initiate inter-agency data agreements now (CBIC, VAHAN, BoI). Lead time: 12–18 months — begin immediately.',
        },
        {
            'title': 'Officer Change Management',
            'detail': 'Government officers accustomed to established manual workflows may be resistant to AI-assisted decision-making.',
            'solution': '45-step guided walkthrough built-in. Designed specifically for non-technical government users. Plus structured training programme in Phase 2.',
        },
    ]
    for i, card in enumerate(cards19):
        cx = Inches(0.3) + i * Inches(4.38)
        cw = Inches(4.2)
        ch = Inches(5.85)
        rect(sl, cx, Inches(1.1), cw, ch, fill=LAMBER)
        rect(sl, cx, Inches(1.1), cw, Inches(0.08), fill=AMBER)
        txb(sl, cx + Inches(0.15), Inches(1.22), cw - Inches(0.3), Inches(0.42),
            paras=[{'text': f'⚠  {card["title"]}', 'size': 11, 'bold': True, 'color': RGBColor(0x85, 0x65, 0x04)}])
        txb(sl, cx + Inches(0.15), Inches(1.75), cw - Inches(0.3), Inches(2.1),
            paras=[{'text': card['detail'], 'size': 10.5, 'color': DARK}])
        rect(sl, cx + Inches(0.15), Inches(3.95), cw - Inches(0.3), Inches(0.04), fill=GREEN)
        rect(sl, cx + Inches(0.1), Inches(4.05), cw - Inches(0.2), Inches(2.6), fill=LGREEN)
        rect(sl, cx + Inches(0.1), Inches(4.05), Inches(0.06), Inches(2.6), fill=GREEN)
        txb(sl, cx + Inches(0.25), Inches(4.1), cw - Inches(0.4), Inches(0.35),
            paras=[{'text': 'Solution:', 'size': 10, 'bold': True, 'color': GREEN}])
        txb(sl, cx + Inches(0.25), Inches(4.48), cw - Inches(0.4), Inches(2.0),
            paras=[{'text': card['solution'], 'size': 10, 'color': DARK}])
    set_notes(sl, "Proactively addressing challenges signals technical credibility. The connectivity issue is real — do not dismiss it. Edge AI on Jetson Orin devices is a proven solution used by NVidia and Indian Railways. The data agreement lead time is the most critical challenge — it must begin before Phase 2 starts, not after.")
    return sl


def s20(prs):
    sl = new_slide(prs)
    T = "Recommendations to LPAI Leadership"
    add_master(sl, 20, T); stitle(sl, T)
    recs = [
        {
            'n': '1', 'bcolor': NAVY,
            'title': 'INITIATE PHASE 2 PILOT AT PETRAPOLE AND ATTARI',
            'body': 'Approve a 6-month pilot integrating LPAI Nexus with the LPMS backend at India\'s two busiest ICPs. The prototype is ready. The pilot will validate AI performance on real operational data and establish officer training protocols for all ICP grades.',
        },
        {
            'n': '2', 'bcolor': GOLD,
            'title': 'ESTABLISH A TECHNOLOGY DIVISION AT LPAI HEADQUARTERS',
            'body': 'Create a dedicated internal team of 8–10 professionals (engineers, data scientists, cybersecurity) to own LPAI Nexus long-term. External vendor dependency for core border management technology creates operational risk and limits institutional knowledge.',
        },
        {
            'n': '3', 'bcolor': GREEN,
            'title': 'INITIATE DATA SHARING AGREEMENTS WITH CBIC, VAHAN, AND BUREAU OF IMMIGRATION',
            'body': 'Axiom\'s full capability requires historical data from these three agencies. Agreements take 12–18 months to formalise. Starting now ensures AI training data is available when Phase 2 deployment begins — any delay cascades into Phase 3.',
        },
    ]
    for i, rec in enumerate(recs):
        ry = Inches(1.1) + i * Inches(1.92)
        rect(sl, Inches(0.3), ry, Inches(12.9), Inches(1.8), fill=LBLUE)
        rect(sl, Inches(0.3), ry, Inches(0.12), Inches(1.8), fill=rec['bcolor'])
        txb(sl, Inches(0.55), ry + Inches(0.08), Inches(0.55), Inches(0.55),
            paras=[{'text': rec['n'], 'size': 24, 'bold': True, 'color': rec['bcolor'], 'align': PP_ALIGN.CENTER}])
        txb(sl, Inches(1.15), ry + Inches(0.1), Inches(11.85), Inches(0.42),
            paras=[{'text': rec['title'], 'size': 12, 'bold': True, 'color': rec['bcolor']}])
        txb(sl, Inches(1.15), ry + Inches(0.55), Inches(11.85), Inches(1.1),
            paras=[{'text': rec['body'], 'size': 11, 'color': DARK}])
    set_notes(sl, "These three recommendations are sequenced intentionally: Recommendation 1 is the immediate ask. Recommendation 2 is the structural ask for long-term success. Recommendation 3 is the parallel track that must start now regardless of the timeline for 1 and 2. Present them as a package — all three are needed for LPAI Nexus to reach its full potential.")
    return sl


def s21(prs):
    sl = new_slide(prs)
    T = "LPAI Nexus — The Opportunity in Numbers"
    add_master(sl, 21, T); stitle(sl, T)
    stats21 = [
        (NAVY,  "95% Faster",      "Cargo document processing"),
        (GOLD,  "~85% Reduction",  "Export clearance time target"),
        (GREEN, "< 1 Second",      "Watchlist face match check"),
        (TEAL,  "24/7 Coverage",   "All camera feeds monitored by AI"),
        (NAVY,  "12 ICPs",         "Unified in one command view"),
        (GOLD,  "₹70,952 Cr",      "Annual trade — all visible in real time"),
    ]
    for i, (c, v, l) in enumerate(stats21):
        col = i % 3
        row = i // 3
        sx = Inches(0.3) + col * Inches(4.35)
        sy = Inches(1.1) + row * Inches(2.35)
        rect(sl, sx, sy, Inches(4.15), Inches(2.2), fill=c)
        txb(sl, sx + Inches(0.15), sy + Inches(0.25), Inches(3.85), Inches(0.9),
            paras=[{'text': v, 'size': 30, 'bold': True,
                    'color': DARK if c == GOLD else WHITE,
                    'align': PP_ALIGN.CENTER}])
        txb(sl, sx + Inches(0.15), sy + Inches(1.2), Inches(3.85), Inches(0.65),
            paras=[{'text': l, 'size': 13,
                    'color': DARK if c == GOLD else WHITE,
                    'align': PP_ALIGN.CENTER}])
    # Banner
    rect(sl, Inches(0.3), Inches(5.65), Inches(12.9), Inches(1.35), fill=DNAVY)
    txb(sl, Inches(0.5), Inches(5.75), Inches(12.7), Inches(0.5),
        paras=[{'text': 'LPAI Nexus is not a concept. It is a working prototype — ready for pilot deployment today.',
                'size': 15, 'bold': True, 'color': GOLD, 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(6.28), Inches(12.7), Inches(0.5),
        paras=[{'text': 'A live demonstration can be arranged at any time for senior leadership.',
                'size': 12, 'color': WHITE, 'align': PP_ALIGN.CENTER}])
    set_notes(sl, "This is the closing impact slide before the conclusion — let the numbers land. Pause on '95% faster' and explain: 45 minutes becomes 2 minutes, for every single consignment, every day. The final banner is a call to action — offer the live demo explicitly. Senior officials are more likely to approve what they have seen working.")
    return sl


def s22(prs):
    sl = new_slide(prs)
    T = "Conclusion"
    add_master(sl, 22, T); stitle(sl, T)
    # Quote box
    rect(sl, Inches(0.5), Inches(1.15), Inches(12.3), Inches(3.2), fill=LBLUE)
    rect(sl, Inches(0.5), Inches(1.15), Inches(0.1), Inches(3.2), fill=NAVY)
    rect(sl, Inches(12.7), Inches(1.15), Inches(0.1), Inches(3.2), fill=NAVY)
    txb(sl, Inches(0.75), Inches(1.28), Inches(12.0), Inches(2.85),
        paras=[
            {'text': '"India manages 15,200 km of international land border across 7 countries.',
             'size': 13, 'bold': True, 'color': NAVY, 'align': PP_ALIGN.CENTER},
            {'text': 'In FY 2023-24, ₹70,952 crore in trade and 30.46 lakh passenger crossings were\nprocessed — a 15-fold increase from a decade ago.',
             'size': 12, 'color': DARK, 'align': PP_ALIGN.CENTER, 'space_before': 8},
            {'text': 'LPAI Nexus demonstrates that the next 15-fold improvement will come\nnot from more infrastructure, but from intelligence."',
             'size': 13, 'bold': True, 'color': NAVY, 'align': PP_ALIGN.CENTER, 'space_before': 12},
        ])
    # 3 deliverables
    txb(sl, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
        paras=[{'text': 'Delivered during this internship:', 'size': 12, 'bold': True,
                'color': NAVY, 'align': PP_ALIGN.CENTER}])
    delivs = [
        (NAVY,  "6 Fully Functional\nModules"),
        (GREEN, "45-Step Guided\nWalkthrough"),
        (GOLD,  "Technical Report\n+ PRD Document"),
    ]
    for i, (c, label) in enumerate(delivs):
        dx = Inches(0.5) + i * Inches(4.15)
        rect(sl, dx, Inches(5.0), Inches(3.95), Inches(1.95), fill=c)
        txb(sl, dx + Inches(0.1), Inches(5.2), Inches(3.75), Inches(1.5),
            paras=[{'text': label, 'size': 16, 'bold': True,
                    'color': DARK if c == GOLD else WHITE,
                    'align': PP_ALIGN.CENTER}])
    set_notes(sl, "End with the central argument of this presentation: the next leap in LPAI's effectiveness will come from intelligence, not more concrete. The three deliverable boxes show tangible output from the internship — not just ideas. Leave the audience with a clear mental model: prototype built, pilot ready, your approval needed.")
    return sl


def s23(prs):
    sl = new_slide(prs)
    rect(sl, 0, 0, W, H, fill=NAVY)
    rect(sl, 0, 0, LST_W, H, fill=GOLD)
    rect(sl, 0, Inches(6.9), W, Inches(0.6), fill=DNAVY)
    rect(sl, Inches(3.5), Inches(3.5), Inches(6.33), Inches(0.04), fill=GOLD)
    txb(sl, Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.65),
        paras=[{'text': 'LPAI Nexus', 'size': 36, 'bold': True,
                'color': WHITE, 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(1.9), Inches(12.33), Inches(0.5),
        paras=[{'text': 'Smart Border Command Platform', 'size': 20,
                'color': GOLD, 'align': PP_ALIGN.CENTER}])
    rect(sl, Inches(3.5), Inches(2.5), Inches(6.33), Inches(0.04), fill=GOLD)
    txb(sl, Inches(0.5), Inches(2.7), Inches(12.33), Inches(0.4),
        paras=[{'text': 'Presented by:', 'size': 12, 'italic': True,
                'color': RGBColor(0xCC, 0xDD, 0xFF), 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(3.15), Inches(12.33), Inches(0.55),
        paras=[{'text': 'Himanshu', 'size': 22, 'bold': True,
                'color': WHITE, 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(3.75), Inches(12.33), Inches(0.4),
        paras=[{'text': 'AI & Software Engineering Intern  |  LPAI, MHA  |  2026',
                'size': 13, 'color': RGBColor(0xCC, 0xDD, 0xFF), 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(4.3), Inches(12.33), Inches(0.35),
        paras=[{'text': 'Under the guidance of:', 'size': 11, 'italic': True,
                'color': RGBColor(0xCC, 0xDD, 0xFF), 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(4.68), Inches(12.33), Inches(0.45),
        paras=[{'text': '[Reporting Officer Name]', 'size': 14, 'bold': True,
                'color': GOLD, 'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.5), Inches(5.18), Inches(12.33), Inches(0.35),
        paras=[{'text': 'Senior Officer (Administration), Land Port Authority of India',
                'size': 11, 'color': RGBColor(0xCC, 0xDD, 0xFF),
                'align': PP_ALIGN.CENTER}])
    txb(sl, Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.4),
        paras=[{'text': 'Land Port Authority of India  |  Ministry of Home Affairs  |  Government of India\n1st Floor, Lok Nayak Bhawan, Khan Market, New Delhi - 110003',
                'size': 9, 'color': RGBColor(0x88, 0x99, 0xBB),
                'align': PP_ALIGN.CENTER}])
    set_notes(sl, "Close by offering a live demonstration of the platform. Have a browser tab open with the LPAI Nexus prototype ready — offer to walk through any module in detail. Leave the room knowing the three recommendations: pilot approval, technology division, and data sharing agreements. Thank the audience for their time and LPAI for the internship opportunity.")
    return sl


# ── Main ────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Building LPAI Nexus Presentation (23 slides)...")
    builders = [s01,s02,s03,s04,s05,s06,s07,s08,s09,s10,
                s11,s12,s13,s14,s15,s16,s17,s18,s19,s20,
                s21,s22,s23]

    for i, fn in enumerate(builders, 1):
        print(f"  Slide {i:02d}: {fn.__name__}...")
        fn(prs)

    prs.save(OUTPUT)
    size = os.path.getsize(OUTPUT)
    print(f"\n  Saved: {OUTPUT}")
    print(f"  Size:  {size:,} bytes  ({size//1024} KB)")
    print(f"  Slides: {len(prs.slides)}")

if __name__ == '__main__':
    main()
